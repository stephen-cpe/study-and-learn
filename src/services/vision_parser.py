import hashlib
import logging
import os
from typing import List, Optional, Tuple

from PIL import Image

from src import db
from src.models import ContentRegistry
from src.services.ai_client import call_ollama
from src.services.exceptions import AIModelUnavailableError
from src.services.vector_store import get_collection_name

logger = logging.getLogger(__name__)


# Vision model default. Migrated from the deprecated qwen3-vl:235b-cloud
# to qwen3.5:397b-cloud (same vendor, Text+Image input, 256K context, Medium
# cloud cost tier). If the model is unreachable on Ollama Cloud, callers
# should expect a soft failure (empty string) + a WARNING log from
# :func:`probe_vision_model_availability` instead of a hard error.
_DEFAULT_VISION_MODEL = "qwen3.5:397b-cloud"
_DEPRECATED_VISION_MODEL = "qwen3-vl:235b-cloud"

# Tracks whether the warning has been logged once per process to avoid
# log spam when the vision model is called repeatedly.
_vision_availability_warned = ""


def probe_vision_model_availability(model: Optional[str] = None) -> bool:
    """Best-effort startup probe for the configured vision model.

    Logs a WARNING and returns ``False`` if the model appears unavailable
    on the active Ollama backend. Returns ``True`` otherwise. Never raises.

    Skipped entirely when ``AI_MOCK=true`` (CI/tests). The probe is lazy and
    idempotent — a WARN is logged at most once per process per model name
    to avoid log spam on repeated calls.

    This is a soft check: it does not block app startup. Real vision
    requests will still attempt the call and degrade gracefully via the
    existing try/except wrappers in :func:`ocr_page` and
    :func:`describe_figure`.
    """
    global _vision_availability_warned

    if os.environ.get("AI_MOCK", "").lower() == "true":
        return True

    if model is None:
        model = os.environ.get("OLLAMA_VISION_MODEL", _DEFAULT_VISION_MODEL)

    # Detect the previously-deprecated default so we can give a useful hint.
    is_deprecated = model == _DEPRECATED_VISION_MODEL

    probe_key = f"{model}:{is_deprecated}"
    if _vision_availability_warned.startswith(f"{probe_key}:"):
        return not _vision_availability_warned.endswith(":unavailable")
    _vision_availability_warned = f"{probe_key}:probing"

    try:
        # Minimal prompt — we only care that the model is reachable and
        # accepts an image-bearing request without returning a "model not
        # found" style error.
        call_ollama(
            "Image: probe\n\nRespond with the single word: ok",
            model=model,
        )
        _vision_availability_warned = f"{probe_key}:available"
        return True
    except AIModelUnavailableError as e:
        hint = (
            f"The configured vision model '{model}' is not reachable on "
            f"the active Ollama backend. Figure descriptions will degrade "
            f"to empty strings. If you are running locally, ensure the "
            f"model is pulled (`ollama pull {model}`). If you are using "
            f"Ollama Cloud, confirm `AI_BACKEND=cloud` is set and your "
            f"API key is valid."
        )
        if is_deprecated:
            hint += (
                f" NOTE: '{model}' is the deprecated default — "
                f"update OLLAMA_VISION_MODEL to "
                f"'{_DEFAULT_VISION_MODEL}' (Qwen3.5:397b-cloud)."
            )
        logger.warning("%s Underlying error: %s", hint, str(e))
        _vision_availability_warned = f"{probe_key}:unavailable"
        return False
    except Exception as e:
        # Non-fatal: the probe is best-effort. Real calls will surface
        # the real error and fall back to empty string.
        logger.debug(
            "Vision model probe returned non-fatal error for '%s': %s",
            model, str(e),
        )
        _vision_availability_warned = f"{probe_key}:available"
        return True


def hash_file(file_path: str) -> str:
    sha256 = hashlib.sha256()
    with open(file_path, 'rb') as f:
        while True:
            chunk = f.read(8192)
            if not chunk:
                break
            sha256.update(chunk)
    return sha256.hexdigest()


def is_content_registered(file_hash: str) -> Optional[str]:
    existing = ContentRegistry.query.filter_by(file_hash=file_hash).first()
    if existing:
        return existing.chroma_collection
    return None


def register_content(file_hash: str, extracted_text: str, ocr_text: str = "") -> str:
    collection_name = get_collection_name(file_hash)
    existing = ContentRegistry.query.filter_by(file_hash=file_hash).first()
    if existing:
        existing.extracted_text = extracted_text
        if ocr_text:
            existing.ocr_text = ocr_text
        db.session.commit()
        return existing.chroma_collection

    entry = ContentRegistry(
        file_hash=file_hash,
        chroma_collection=collection_name,
        extracted_text=extracted_text,
        ocr_text=ocr_text,
    )
    db.session.add(entry)
    try:
        db.session.commit()
    except Exception:
        db.session.rollback()
        existing = ContentRegistry.query.filter_by(file_hash=file_hash).first()
        if existing:
            return existing.chroma_collection
        raise
    return collection_name


def _get_poppler_path() -> Optional[str]:
    env_poppler = os.environ.get("POPPLER_PATH", "")
    if env_poppler and os.path.isdir(env_poppler):
        return env_poppler

    possible = [
        r"C:\Program Files\poppler\Library\bin",
        r"C:\Program Files\poppler\bin",
        r"C:\poppler\Library\bin",
        r"C:\poppler\bin",
    ]
    for p in possible:
        if os.path.isdir(p):
            return p
    return None


def render_pdf_page(file_path: str, page_number: int, output_dir: str) -> str:
    poppler_path = _get_poppler_path()
    kwargs = {}
    if poppler_path:
        kwargs["poppler_path"] = poppler_path

    from pdf2image import convert_from_path
    images = convert_from_path(file_path, first_page=page_number + 1,
                               last_page=page_number + 1, dpi=200, **kwargs)
    if not images:
        raise RuntimeError(f"Failed to render page {page_number + 1} of {file_path}")

    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, f"page_{page_number + 1}.png")
    images[0].save(output_path, "PNG")
    return output_path


def render_pdf_pages(file_path: str, output_dir: str) -> List[str]:
    poppler_path = _get_poppler_path()
    kwargs = {}
    if poppler_path:
        kwargs["poppler_path"] = poppler_path

    from pdf2image import convert_from_path
    from pdf2image.exceptions import PDFInfoNotInstalledError
    try:
        images = convert_from_path(file_path, dpi=200, **kwargs)
    except PDFInfoNotInstalledError:
        logger.warning("Poppler not found — cannot render PDF pages. Install poppler-utils.")
        return []

    os.makedirs(output_dir, exist_ok=True)
    paths = []
    for i, img in enumerate(images):
        out_path = os.path.join(output_dir, f"page_{i + 1}.png")
        img.save(out_path, "PNG")
        paths.append(out_path)
    return paths


def extract_docx_images(file_path: str, output_dir: str) -> List[str]:
    try:
        from docx import Document
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
    except ImportError:
        logger.warning("python-docx not available for image extraction")
        return []

    doc = Document(file_path)
    os.makedirs(output_dir, exist_ok=True)
    paths = []

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image = rel.target_part
            ext = os.path.splitext(image.partname)[-1] or ".png"
            out_path = os.path.join(output_dir, f"docx_image_{image.partname.replace('/', '_')}{ext}")
            with open(out_path, "wb") as f:
                f.write(image.blob)
            try:
                img = Image.open(out_path)
                if img.mode != "RGB":
                    img = img.convert("RGB")
                png_path = os.path.splitext(out_path)[0] + ".png"
                img.save(png_path, "PNG")
                if png_path != out_path:
                    paths.append(png_path)
                else:
                    paths.append(out_path)
            except Exception:
                paths.append(out_path)

    return paths


def extract_pptx_content(file_path: str, output_dir: str) -> Tuple[str, List[str]]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        logger.warning("python-pptx not available")
        return "", []

    prs = Presentation(file_path)
    text_parts = []
    os.makedirs(output_dir, exist_ok=True)
    slide_images = []

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        text_parts.append(paragraph.text)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    ext = image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    img_path = os.path.join(
                        output_dir, f"slide_{i + 1}_img_{len(slide_images)}.{ext}"
                    )
                    with open(img_path, "wb") as f:
                        f.write(image.blob)
                    try:
                        img = Image.open(img_path)
                        if img.mode != "RGB":
                            img = img.convert("RGB")
                        png_path = os.path.splitext(img_path)[0] + ".png"
                        img.save(png_path, "PNG")
                        if png_path != img_path:
                            os.remove(img_path)
                        slide_images.append(png_path)
                    except Exception:
                        slide_images.append(img_path)
                except Exception:
                    pass

    return "\n".join(text_parts), slide_images


def _resize_image_if_needed(image_path: str) -> str:
    max_dim = int(os.environ.get("OCR_MAX_IMAGE_DIMENSION", "2048"))
    try:
        img = Image.open(image_path)
        w, h = img.size
        if w <= max_dim and h <= max_dim:
            return image_path
        ratio = min(max_dim / w, max_dim / h)
        new_size = (int(w * ratio), int(h * ratio))
        img = img.resize(new_size, Image.LANCZOS)
        img.save(image_path)
    except Exception as e:
        logger.warning("Failed to resize image %s: %s", image_path, str(e))
    return image_path


def ocr_page(image_path: str, mode: str = "text") -> str:
    max_bytes = int(os.environ.get("OCR_MAX_FILE_BYTES", str(50 * 1024 * 1024)))
    file_size = os.path.getsize(image_path)
    if file_size > max_bytes:
        logger.warning(
            "Skipping OCR for %s: file size %d exceeds OCR_MAX_FILE_BYTES (%d)",
            image_path, file_size, max_bytes,
        )
        return ""
    _resize_image_if_needed(image_path)
    abs_path = os.path.abspath(image_path)

    mode_prompts = {
        "text": f"Text Recognition: {abs_path}",
        "table": f"Table Recognition: {abs_path}",
        "figure": f"Figure Recognition: {abs_path}",
    }
    if mode not in mode_prompts:
        raise ValueError(f"Unknown OCR mode: {mode}. Use 'text', 'table', or 'figure'.")

    prompt = mode_prompts[mode]
    model = os.environ.get("OLLAMA_OCR_MODEL", "glm-ocr")
    timeout = int(os.environ.get("OCR_TIMEOUT_PER_PAGE", "120"))

    try:
        result = call_ollama(prompt, model=model, force_local=True)
        return result.strip() if result else ""
    except Exception as e:
        logger.warning("OCR %s mode failed for %s: %s", mode, image_path, str(e))
        return ""


def ocr_page_full(image_path: str) -> str:
    run_full = os.environ.get("OCR_FULL", "").lower() == "true"
    modes = ["text", "table", "figure"] if run_full else ["text"]
    parts = []
    for mode in modes:
        try:
            result = ocr_page(image_path, mode=mode)
            if result and result.strip():
                label = {"text": "Text", "table": "Table", "figure": "Figure"}[mode]
                parts.append(f"[{label} OCR]\n{result}")
        except Exception as e:
            logger.warning("OCR %s mode error on %s: %s", mode, image_path, str(e))

    return "\n\n".join(parts)


def describe_figure(image_path: str) -> str:
    skip = os.environ.get("OCR_FIGURE_DESCRIPTION", "false").lower() != "true"
    if skip:
        return ""

    _resize_image_if_needed(image_path)
    abs_path = os.path.abspath(image_path)

    model = os.environ.get("OLLAMA_VISION_MODEL", _DEFAULT_VISION_MODEL)

    # Best-effort availability check. Logs a WARNING once per process if
    # the configured model is unreachable on Ollama Cloud. Never raises.
    probe_vision_model_availability(model)

    prompt = (
        f"Image: {abs_path}\n\n"
        "Describe what this figure explains in 2-3 sentences, focusing on the key concepts "
        "and how they relate to each other. Include any labels, axis titles, or annotations "
        "visible in the figure."
    )

    try:
        result = call_ollama(prompt, model=model)
        return result.strip() if result else ""
    except Exception as e:
        logger.warning("Figure description failed for %s: %s", image_path, str(e))
        return ""


def extract_text_with_vision(file_path: str, progress_callback=None) -> str:
    from src.services.document_parser import extract_text as _basic_extract

    ext = os.path.splitext(file_path)[1].lower()

    file_hash = hash_file(file_path)
    existing_collection = is_content_registered(file_hash)
    if existing_collection:
        entry = ContentRegistry.query.filter_by(file_hash=file_hash).first()
        if entry and entry.extracted_text:
            logger.info("Content already registered for hash %s, returning cached text", file_hash)
            return entry.extracted_text

    parts = []

    try:
        basic_text = _basic_extract(file_path)
        if basic_text and basic_text.strip():
            parts.append(basic_text)
    except Exception as e:
        logger.warning("Basic text extraction failed for %s: %s", file_path, str(e))

    if ext in ('.txt', '.md'):
        result = basic_text if parts else ""
        register_content(file_hash, result)
        return result

    ocr_enabled = os.environ.get("OCR_FULL", "").lower() == "true"

    if not ocr_enabled and ext in ('.pdf', '.docx', '.pptx'):
        result = "\n\n".join(parts) if parts else ""
        if not result or not result.strip():
            result = ""
        register_content(file_hash, result)
        return result

    output_dir = os.path.join(os.path.dirname(file_path), "_ocr_temp")
    os.makedirs(output_dir, exist_ok=True)

    total_pages = 0
    page_images = []

    if ext == '.pdf':
        page_images = render_pdf_pages(file_path, output_dir)
        total_pages = len(page_images)
    elif ext == '.docx':
        docx_imgs = extract_docx_images(file_path, output_dir)
        page_images = docx_imgs
        total_pages = len(page_images)
    elif ext == '.pptx':
        pptx_text, slide_imgs = extract_pptx_content(file_path, output_dir)
        if pptx_text and pptx_text.strip():
            parts.append(pptx_text)
        page_images = slide_imgs
        total_pages = len(page_images)
    elif ext in ('.png', '.jpg', '.jpeg'):
        page_images = [file_path]
        total_pages = 1

    ocr_outputs = []
    figure_outputs = []

    for i, img_path in enumerate(page_images):
        try:
            if progress_callback:
                progress_callback("ocr", i + 1, total_pages)

            ocr_result = ocr_page_full(img_path)
            if ocr_result:
                ocr_outputs.append(ocr_result)
        except Exception as e:
            logger.warning("OCR failed for %s: %s", img_path, str(e))

        try:
            if progress_callback:
                progress_callback("figure", i + 1, total_pages)

            fig_desc = describe_figure(img_path)
            if fig_desc:
                figure_outputs.append(f"[Figure Description]\n{fig_desc}")
        except Exception as e:
            logger.warning("Figure description failed for %s: %s", img_path, str(e))

    if ocr_outputs:
        parts.append("[BEGIN OCR OUTPUT]\n" + "\n\n".join(ocr_outputs) + "\n[END OCR OUTPUT]")
    if figure_outputs:
        parts.append("[BEGIN FIGURE DESCRIPTIONS]\n" + "\n\n".join(figure_outputs) + "\n[END FIGURE DESCRIPTIONS]")

    result = "\n\n".join(parts)
    if not result or not result.strip():
        if parts:
            result = parts[0]
        else:
            result = ""

    ocr_combined = "\n\n".join(ocr_outputs) if ocr_outputs else ""
    register_content(file_hash, result, ocr_text=ocr_combined)

    for img_path in page_images:
        try:
            if os.path.exists(img_path) and img_path != file_path:
                os.remove(img_path)
        except Exception:
            pass
    try:
        if os.path.exists(output_dir) and output_dir != os.path.dirname(file_path):
            remaining = os.listdir(output_dir)
            if not remaining:
                os.rmdir(output_dir)
    except Exception:
        pass

    return result
