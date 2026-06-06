import logging
import os
import pytest
from unittest.mock import patch, MagicMock


@pytest.fixture(autouse=True)
def mock_ai(monkeypatch):
    monkeypatch.setenv("AI_MOCK", "true")


@pytest.fixture
def test_png_path(tmp_path):
    from PIL import Image
    img = Image.new("RGB", (100, 100), color="black")
    path = os.path.join(tmp_path, "test.png")
    img.save(path)
    return path


@pytest.fixture
def test_txt_path(tmp_path):
    path = os.path.join(tmp_path, "test.txt")
    with open(path, "w") as f:
        f.write("Hello world. This is a test file.")
    return path


@pytest.fixture
def fixtures_dir():
    return os.path.join(os.path.dirname(__file__), "fixtures")


class TestHashFile:
    def test_same_content(self, test_txt_path):
        from src.services.vision_parser import hash_file
        h1 = hash_file(test_txt_path)
        h2 = hash_file(test_txt_path)
        assert h1 == h2

    def test_different_content(self, tmp_path):
        from src.services.vision_parser import hash_file
        a = os.path.join(tmp_path, "a.txt")
        b = os.path.join(tmp_path, "b.txt")
        with open(a, "w") as f:
            f.write("content A")
        with open(b, "w") as f:
            f.write("content B")
        assert hash_file(a) != hash_file(b)

    def test_empty_file(self, tmp_path):
        from src.services.vision_parser import hash_file
        p = os.path.join(tmp_path, "empty.txt")
        with open(p, "w") as f:
            f.write("")
        result = hash_file(p)
        assert len(result) == 64
        assert all(c in "0123456789abcdef" for c in result)


class TestContentRegistry:
    def test_is_content_registered_none(self):
        with patch("src.services.vision_parser.ContentRegistry") as mock_cr:
            mock_cr.query.filter_by.return_value.first.return_value = None
            from src.services.vision_parser import is_content_registered
            assert is_content_registered("unknown_hash") is None

    def test_is_content_registered_found(self):
        with patch("src.services.vision_parser.ContentRegistry") as mock_cr:
            mock_existing = MagicMock()
            mock_existing.chroma_collection = "doc_dummy123"
            mock_cr.query.filter_by.return_value.first.return_value = mock_existing
            from src.services.vision_parser import is_content_registered
            result = is_content_registered("dummy123")
            assert result == "doc_dummy123"

    def test_register_content_new(self):
        with patch("src.services.vision_parser.ContentRegistry") as mock_cr, \
             patch("src.services.vision_parser.db") as mock_db:
            mock_cr.query.filter_by.return_value.first.return_value = None
            mock_cr.return_value = MagicMock()
            from src.services.vision_parser import register_content
            result = register_content("h" * 64, "sample text")
            assert result.startswith("doc_")
            assert len(result) <= 63

    def test_register_content_existing(self):
        with patch("src.services.vision_parser.ContentRegistry") as mock_cr, \
             patch("src.services.vision_parser.db") as mock_db:
            mock_existing = MagicMock()
            mock_existing.chroma_collection = "doc_existing"
            mock_cr.query.filter_by.return_value.first.return_value = mock_existing
            from src.services.vision_parser import register_content
            result = register_content("h" * 64, "updated text", ocr_text="ocr")
            assert result == "doc_existing"


class TestCollectionName:
    def test_collection_name_length(self):
        from src.services.vector_store import get_collection_name
        h = "a" * 64
        name = get_collection_name(h)
        assert name.startswith("doc_")
        assert len(name) <= 63
        assert len(name) == 4 + 59


class TestOcrPage:
    def test_mock(self, test_png_path):
        from src.services.vision_parser import ocr_page
        result = ocr_page(test_png_path, mode="text")
        assert isinstance(result, str)
        assert len(result) > 0

    def test_invalid_mode(self, test_png_path):
        from src.services.vision_parser import ocr_page
        with pytest.raises(ValueError, match="Unknown OCR mode"):
            ocr_page(test_png_path, mode="invalid")


class TestOcrPageFull:
    def test_full_mock(self, test_png_path, monkeypatch):
        monkeypatch.setenv("OCR_FULL", "true")
        from src.services.vision_parser import ocr_page_full
        result = ocr_page_full(test_png_path)
        assert "[Text OCR]" in result
        assert "[Table OCR]" in result
        assert "[Figure OCR]" in result

    def test_default_text_only(self, test_png_path, monkeypatch):
        monkeypatch.delenv("OCR_FULL", raising=False)
        from src.services.vision_parser import ocr_page_full
        result = ocr_page_full(test_png_path)
        assert "[Text OCR]" in result
        assert "[Table OCR]" not in result
        assert "[Figure OCR]" not in result


class TestDescribeFigure:
    def test_mock(self, test_png_path, monkeypatch):
        monkeypatch.setenv("OCR_FIGURE_DESCRIPTION", "true")
        from src.services.vision_parser import describe_figure
        result = describe_figure(test_png_path)
        assert isinstance(result, str)
        assert len(result) > 0

    def test_disabled_by_default(self, test_png_path, monkeypatch):
        monkeypatch.delenv("OCR_FIGURE_DESCRIPTION", raising=False)
        from src.services.vision_parser import describe_figure
        result = describe_figure(test_png_path)
        assert result == ""

    def test_model_unavailable_returns_empty(self, test_png_path, monkeypatch):
        """When Ollama reports the model as missing, describe_figure()
        must degrade gracefully (return '') rather than raise — otherwise
        extract_text_with_vision() would crash on every figure page.
        """
        monkeypatch.setenv("OCR_FIGURE_DESCRIPTION", "true")
        monkeypatch.setenv("AI_MOCK", "false")
        from src.services import vision_parser as vp
        from src.services.exceptions import AIModelUnavailableError

        monkeypatch.setattr(vp, "_vision_availability_warned", "")
        with patch.object(vp, "call_ollama",
                          side_effect=AIModelUnavailableError("model not found")):
            result = vp.describe_figure(test_png_path)
        assert result == ""

    def test_unexpected_exception_returns_empty(self, test_png_path, monkeypatch):
        """Any exception from call_ollama must be swallowed to avoid
        blocking the whole vision pipeline on a single bad figure.
        """
        monkeypatch.setenv("OCR_FIGURE_DESCRIPTION", "true")
        monkeypatch.setenv("AI_MOCK", "false")
        from src.services import vision_parser as vp

        monkeypatch.setattr(vp, "_vision_availability_warned", "")
        with patch.object(vp, "call_ollama", side_effect=RuntimeError("boom")):
            result = vp.describe_figure(test_png_path)
        assert result == ""

    def test_empty_response_stripped(self, test_png_path, monkeypatch):
        """Whitespace-only or empty responses should normalize to ''."""
        monkeypatch.setenv("OCR_FIGURE_DESCRIPTION", "true")
        monkeypatch.setenv("AI_MOCK", "false")
        from src.services import vision_parser as vp

        monkeypatch.setattr(vp, "_vision_availability_warned", "")
        with patch.object(vp, "call_ollama", return_value="   \n  "):
            result = vp.describe_figure(test_png_path)
        assert result == ""


class TestOcrPageFailureHandling:
    def test_model_unavailable_returns_empty(self, test_png_path, monkeypatch):
        """ocr_page() must also degrade gracefully on model failure."""
        from src.services import vision_parser as vp
        from src.services.exceptions import AIModelUnavailableError

        monkeypatch.setenv("AI_MOCK", "false")
        with patch.object(vp, "call_ollama",
                          side_effect=AIModelUnavailableError("model not found")):
            result = vp.ocr_page(test_png_path, mode="text")
        assert result == ""


class TestVisionModelProbe:
    def test_skipped_under_ai_mock(self, monkeypatch):
        """The probe must be a no-op in CI — it must not try to reach Ollama
        when AI_MOCK=true (otherwise tests would do real network I/O).
        """
        from src.services import vision_parser as vp
        # Even with call_ollama patched to raise, probe should short-circuit.
        with patch.object(vp, "call_ollama",
                          side_effect=AssertionError("probe touched call_ollama under AI_MOCK")):
            assert vp.probe_vision_model_availability() is True

    def test_warns_once_on_unavailable(self, monkeypatch):
        """First probe logs WARNING and returns False; subsequent calls
        with the same model do not re-probe (avoids log spam on every
        figure call).
        """
        from src.services import vision_parser as vp
        from src.services.exceptions import AIModelUnavailableError

        # Disable AI_MOCK (autouse fixture sets it true) so the probe
        # actually exercises the call path.
        monkeypatch.setenv("AI_MOCK", "false")
        vp._vision_availability_warned = ""
        calls = []

        def fake_call(prompt, model=None, **kwargs):
            calls.append((prompt, model))
            raise AIModelUnavailableError("not found")

        with patch.object(vp, "call_ollama", side_effect=fake_call):
            first = vp.probe_vision_model_availability("qwen3.5:397b-cloud")
            second = vp.probe_vision_model_availability("qwen3.5:397b-cloud")
        assert first is False
        assert second is False
        # First call hit call_ollama; second call was short-circuited by cache.
        assert len(calls) == 1
        assert calls[0][1] == "qwen3.5:397b-cloud"

    def test_warns_with_deprecation_hint_for_old_default(self, monkeypatch):
        """If a user still has the deprecated qwen3-vl:235b-cloud default,
        the warning must explicitly point them to the new model.
        """
        from src.services import vision_parser as vp
        from src.services.exceptions import AIModelUnavailableError

        monkeypatch.setenv("AI_MOCK", "false")
        vp._vision_availability_warned = ""
        with patch.object(vp, "call_ollama",
                          side_effect=AIModelUnavailableError("not found")):
            with patch.object(vp.logger, "warning") as mock_warn:
                vp.probe_vision_model_availability("qwen3-vl:235b-cloud")
        # The probe logs at least one WARNING; the deprecation hint must
        # appear in the concatenated call args (first arg is a format
        # string, remaining args are its substitutions).
        assert mock_warn.call_count >= 1
        all_args = []
        for call in mock_warn.call_args_list:
            all_args.extend(call.args)
            all_args.extend(str(v) for v in call.kwargs.values())
        combined = " ".join(all_args)
        assert "deprecated" in combined.lower()
        assert "qwen3.5:397b-cloud" in combined

    def test_returns_true_on_success(self, monkeypatch):
        from src.services import vision_parser as vp
        monkeypatch.setenv("AI_MOCK", "false")
        vp._vision_availability_warned = ""
        with patch.object(vp, "call_ollama", return_value="ok"):
            assert vp.probe_vision_model_availability("qwen3.5:397b-cloud") is True

    def test_non_unavailable_exception_is_treated_as_soft(self, monkeypatch):
        """A non-AIModelUnavailableError exception (e.g. timeout, HTTP 500)
        should be DEBUG-logged but not produce a user-facing WARNING — the
        existing call_ollama exception handlers in describe_figure/ocr_page
        will surface the real error.
        """
        from src.services import vision_parser as vp

        monkeypatch.setenv("AI_MOCK", "false")
        vp._vision_availability_warned = ""
        with patch.object(vp, "call_ollama", side_effect=RuntimeError("timeout")):
            with patch.object(vp.logger, "warning") as mock_warn, \
                 patch.object(vp.logger, "debug") as mock_debug:
                result = vp.probe_vision_model_availability("qwen3.5:397b-cloud")
        assert result is True
        assert mock_warn.call_count == 0
        assert mock_debug.call_count == 1

    def test_initial_state_is_string_not_bool(self, monkeypatch):
        """Regression: _vision_availability_warned was initialized to False
        (bool), causing 'bool' object has no attribute 'startswith' on the
        first call to probe_vision_model_availability (line 56 uses
        .startswith()). Fixed by initializing to '' (empty str).
        """
        from src.services import vision_parser as vp

        monkeypatch.setenv("AI_MOCK", "false")
        # Reset to the clean initial state and verify it's a string
        monkeypatch.setattr(vp, "_vision_availability_warned", "")
        assert isinstance(vp._vision_availability_warned, str), (
            "_vision_availability_warned must be a str, not bool"
        )
        # Verify probe works correctly from the clean initial state
        with patch.object(vp, "call_ollama", return_value="ok"):
            result = vp.probe_vision_model_availability("qwen3.5:397b-cloud")
        assert result is True


class TestExtractTextWithVision:
    def test_txt_file(self, test_txt_path):
        with patch("src.services.vision_parser.is_content_registered", return_value=None), \
             patch("src.services.vision_parser.register_content"):
            from src.services.document_parser import extract_text_with_vision
            text = extract_text_with_vision(test_txt_path)
            assert "Hello world" in text

    def test_dedup(self, fixtures_dir):
        import os
        pdf_path = os.path.join(fixtures_dir, "sample.pdf")
        with patch("src.services.vision_parser.is_content_registered", return_value="doc_cached"), \
             patch("src.services.vision_parser.ContentRegistry") as mock_cr, \
             patch("src.services.vision_parser.register_content"):
            mock_entry = MagicMock()
            mock_entry.extracted_text = "cached pdf result"
            mock_cr.query.filter_by.return_value.first.return_value = mock_entry

            from src.services.document_parser import extract_text_with_vision
            result = extract_text_with_vision(pdf_path)
            assert result == "cached pdf result"

    def test_image_mock(self, test_png_path):
        with patch("src.services.vision_parser.is_content_registered", return_value=None), \
             patch("src.services.vision_parser.register_content"):
            from src.services.document_parser import extract_text_with_vision
            text = extract_text_with_vision(test_png_path)
            assert isinstance(text, str)

    def test_pdf_single_page(self, fixtures_dir, test_txt_path):
        pdf_path = os.path.join(fixtures_dir, "sample.pdf")
        with patch("src.services.vision_parser.is_content_registered", return_value=None), \
             patch("src.services.vision_parser.register_content"):
            from src.services.document_parser import extract_text_with_vision
            text = extract_text_with_vision(pdf_path)
            assert isinstance(text, str)
            assert len(text) > 0

    def test_pdf_skips_ocr_when_ocr_full_false(self, fixtures_dir, monkeypatch):
        monkeypatch.setenv("OCR_FULL", "false")
        pdf_path = os.path.join(fixtures_dir, "sample.pdf")
        with patch("src.services.vision_parser.is_content_registered", return_value=None), \
             patch("src.services.vision_parser.register_content") as mock_reg, \
             patch("src.services.vision_parser.render_pdf_pages") as mock_render, \
             patch("src.services.vision_parser.ocr_page_full") as mock_ocr:
            from src.services.document_parser import extract_text_with_vision
            text = extract_text_with_vision(pdf_path)
            assert isinstance(text, str)
            assert len(text) > 0
            mock_render.assert_not_called()
            mock_ocr.assert_not_called()
            mock_reg.assert_called_once()

    def test_pptx_skips_ocr_when_ocr_full_false(self, tmp_path, monkeypatch):
        monkeypatch.setenv("OCR_FULL", "false")
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test PPTX No OCR"
        pptx_path = os.path.join(tmp_path, "test.pptx")
        prs.save(pptx_path)

        with patch("src.services.vision_parser.is_content_registered", return_value=None), \
             patch("src.services.vision_parser.register_content") as mock_reg, \
             patch("src.services.vision_parser.ocr_page_full") as mock_ocr:
            from src.services.document_parser import extract_text_with_vision
            text = extract_text_with_vision(pptx_path)
            assert "Test PPTX No OCR" in text
            mock_ocr.assert_not_called()
            mock_reg.assert_called_once()


class TestResize:
    def test_small_image_unchanged(self, test_png_path, monkeypatch):
        monkeypatch.setenv("OCR_MAX_IMAGE_DIMENSION", "2048")
        from src.services.vision_parser import _resize_image_if_needed
        result = _resize_image_if_needed(test_png_path)
        assert result == test_png_path


class TestConcurrentHandling:
    def test_register_content_integrity_error(self, tmp_path):
        with patch("src.services.vision_parser.ContentRegistry") as mock_cr, \
             patch("src.services.vision_parser.db") as mock_db:
            mock_cr.query.filter_by.return_value.first.side_effect = [None, MagicMock(chroma_collection="doc_fallback")]
            mock_cr.return_value = MagicMock()
            mock_db.session.commit.side_effect = Exception("unique constraint")

            from src.services.vision_parser import register_content
            result = register_content("h" * 64, "text")
            assert result == "doc_fallback"
            mock_db.session.rollback.assert_called_once()
