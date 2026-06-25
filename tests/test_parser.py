"""
Unit tests for the document parser service.
"""
import os
import tempfile
import pytest
from unittest.mock import patch, MagicMock
from src.services.document_parser import extract_text


def test_extract_valid_txt():
    """Test extracting text from a valid .txt file."""
    with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as f:
        f.write('Hello, world!')
        temp_path = f.name
    
    try:
        text = extract_text(temp_path)
        assert text == 'Hello, world!'
    finally:
        os.unlink(temp_path)


def test_extract_valid_md():
    """Test extracting text from a valid .md file."""
    with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False) as f:
        f.write('# Hello\\n**bold** and _italic_')
        temp_path = f.name
    
    try:
        text = extract_text(temp_path)
        # .md parser returns raw content (trivial markdown stripping is intentionally not performed).
        assert text == '# Hello\\n**bold** and _italic_'
    finally:
        os.unlink(temp_path)


def test_extract_empty_file_raises():
    """Test that extracting text from an empty file raises ValueError."""
    with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as f:
        # Write nothing
        temp_path = f.name
    
    try:
        with pytest.raises(ValueError, match='File is empty'):
            extract_text(temp_path)
    finally:
        os.unlink(temp_path)


def test_extract_nonexistent_path_raises():
    """Test that extracting text from a nonexistent path raises ValueError."""
    with pytest.raises(ValueError, match='File not found'):
        extract_text('/nonexistent/path/file.txt')


def test_extract_text_pptx(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test PPTX Slide"
    pptx_path = os.path.join(tmp_path, "test.pptx")
    prs.save(pptx_path)

    text = extract_text(pptx_path)
    assert "Test PPTX Slide" in text


def test_extract_text_image_returns_empty(tmp_path):
    from PIL import Image
    img_path = os.path.join(tmp_path, "test_image.png")
    img = Image.new("RGB", (10, 10), color="white")
    img.save(img_path)

    text = extract_text(img_path)
    assert text == ''


def test_extract_text_with_vision_txt_delegates(tmp_path):
    from src.services.document_parser import extract_text_with_vision

    txt_path = os.path.join(tmp_path, "test.txt")
    with open(txt_path, "w") as f:
        f.write("Hello world from vision txt test")

    with patch("src.services.vision_parser.is_content_registered", return_value="doc_skip"), \
         patch("src.services.vision_parser.register_content"):
        result = extract_text_with_vision(txt_path)
    assert result == "Hello world from vision txt test"


def test_extract_text_with_vision_md_registers_content(tmp_path):
    from src.services.document_parser import extract_text_with_vision

    md_path = os.path.join(tmp_path, "test.md")
    with open(md_path, "w") as f:
        f.write("# Hello Markdown\n\nSome content here.")

    with patch("src.services.vision_parser.hash_file") as mock_hash, \
         patch("src.services.vision_parser.is_content_registered") as mock_is_reg, \
         patch("src.services.vision_parser.register_content") as mock_reg:
        mock_hash.return_value = "a" * 64
        mock_is_reg.return_value = None

        result = extract_text_with_vision(md_path)

        assert result == "# Hello Markdown\n\nSome content here."
        mock_is_reg.assert_called_once_with("a" * 64)
        mock_reg.assert_called_once_with("a" * 64, "# Hello Markdown\n\nSome content here.")


def test_extract_text_with_vision_txt_skips_registration_if_already_registered(tmp_path):
    from src.services.document_parser import extract_text_with_vision

    txt_path = os.path.join(tmp_path, "test.txt")
    with open(txt_path, "w") as f:
        f.write("Already registered text")

    with patch("src.services.vision_parser.hash_file") as mock_hash, \
         patch("src.services.vision_parser.is_content_registered") as mock_is_reg, \
         patch("src.services.vision_parser.register_content") as mock_reg:
        mock_hash.return_value = "b" * 64
        mock_is_reg.return_value = "doc_existing_coll"

        result = extract_text_with_vision(txt_path)

        assert result == "Already registered text"
        mock_is_reg.assert_called_once_with("b" * 64)
        mock_reg.assert_not_called()


def test_extract_text_with_vision_unknown_ext(tmp_path):
    from src.services.document_parser import extract_text_with_vision

    bad_path = os.path.join(tmp_path, "test.xyz")
    with open(bad_path, "w") as f:
        f.write("nope")

    with patch("src.services.vision_parser.is_content_registered", return_value=None), \
         patch("src.services.vision_parser.register_content"):
        result = extract_text_with_vision(bad_path)
        assert result == ""


def test_extract_text_with_vision_pdf_dedup(tmp_path, monkeypatch):
    monkeypatch.setenv("AI_MOCK", "true")
    from src.services.document_parser import extract_text_with_vision

    import os
    fixtures = os.path.join(os.path.dirname(__file__), "fixtures", "sample.pdf")

    with patch("src.services.vision_parser.is_content_registered") as mock_is_reg, \
         patch("src.services.vision_parser.register_content"), \
         patch("src.services.vision_parser.ContentRegistry"):
        mock_is_reg.return_value = None
        text1 = extract_text_with_vision(fixtures)

        mock_is_reg.return_value = "doc_cached"
        mock_entry = MagicMock()
        mock_entry.extracted_text = text1
        with patch("src.services.vision_parser.ContentRegistry") as mock_cr:
            mock_cr.query.filter_by.return_value.first.return_value = mock_entry
            text2 = extract_text_with_vision(fixtures)

    assert text1 == text2
    assert len(text1) > 0