"""
Document parser service for extracting text from uploaded files.
"""
import os
import logging

logger = logging.getLogger(__name__)


def extract_text(file_path: str) -> str:
    """Extract text from a file based on its extension.
    
    Supports .txt, .md, .pdf, .docx, and .pptx files.
    
    Args:
        file_path (str): Path to the file to extract text from
        
    Returns:
        str: The extracted text content
        
    Raises:
        ValueError: If file is empty, not found, or unsupported type
    """
    if not os.path.exists(file_path):
        raise ValueError('File not found')
    
    if os.path.getsize(file_path) == 0:
        raise ValueError('File is empty')
    
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    if ext == '.txt' or ext == '.md':
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
    elif ext == '.pdf':
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            text = ''
            for page in reader.pages:
                text += page.extract_text() + '\n'
        except Exception as e:
            raise ValueError(f'Failed to extract text from PDF: {str(e)}')
    elif ext == '.docx':
        try:
            from docx import Document
            doc = Document(file_path)
            text = '\n'.join([para.text for para in doc.paragraphs])
        except Exception as e:
            raise ValueError(f'Failed to extract text from DOCX: {str(e)}')
    elif ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                text_parts.append(paragraph.text)
            text = '\n'.join(text_parts)
        except Exception as e:
            raise ValueError(f'Failed to extract text from PPTX: {str(e)}')
    elif ext in ('.png', '.jpg', '.jpeg'):
        return ''
    else:
        raise ValueError(f'Unsupported file type: {ext}')
    
    if not text or not text.strip():
        raise ValueError('No readable text found in file')
    
    return text


def extract_text_with_vision(file_path: str, progress_callback=None) -> str:
    """Enhanced text extraction using OCR and vision models when beneficial.
    
    For .txt and .md files, delegates to the basic extract_text().
    For .pdf, .docx, .pptx, and image files, uses the vision pipeline:
    computes file hash, checks ContentRegistry for cached results,
    and runs OCR/vision analysis only for new content.
    
    Args:
        file_path: Path to the uploaded file
        progress_callback: Optional callable(stage, current, total) for UI updates
    
    Returns:
        Enriched text corpus with OCR and figure descriptions merged in
    """
    from src.services.vision_parser import (
        hash_file, is_content_registered, extract_text_with_vision as _vision_extract
    )

    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    if ext in ('.txt', '.md'):
        return extract_text(file_path)

    return _vision_extract(file_path, progress_callback=progress_callback)